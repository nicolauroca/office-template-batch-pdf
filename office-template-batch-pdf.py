import os
import re
import csv
import sys
import json
import shutil
import atexit
import logging
import argparse
import tempfile
import subprocess
import pandas as pd
from pathlib import Path
from typing import Optional, Dict, Any

# PPTX
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# DOCX
from docx import Document as DocxDocument

try:
    from tqdm import tqdm  # type: ignore
except Exception:
    tqdm = None  # sin barra de progreso

__version__ = "0.9.0"

"""
Batch DOC/DOCX/PPT/PPTX → PDF with {{TOKEN}} substitution from a spreadsheet.

New/extended features:
- Preflight token vs. columns (strict optional).
- Dry-run mode.
- Row-level SKIP and per-row subfolder (OUTPUT).
- Legacy conversion cache (.doc/.ppt/.odt/.odp/.rtf → OOXML).
- Retries on PDF export.
- JSON and CSV reports.
- CSV input supported besides XLSX (auto by extension).
- Progress bar via tqdm if available.
- Token filters and default values: {{Field|trim|upper}}, {{Amount|euros}}, {{Date|dmy}}, {{Field?:N/A}}.
- Additional template formats: .odt, .odp, .rtf (converted to OOXML prior to replacement).
- Pluggable export engine: LibreOffice or Microsoft Office (Word/PowerPoint COM on Windows, if available).

Requirements:
  pip install pandas python-pptx python-docx
  (optional) pip install tqdm pywin32 or comtypes (for MS Office on Windows)
LibreOffice must be installed and 'soffice' available in PATH (or set SOFFICE_BIN) for conversions
to/from OOXML and/or for PDF export when chosen.
"""

# =========================
# ⚙️ CONFIGURATION
# =========================

try:
    _HERE = Path(__file__).resolve().parent
except NameError:
    _HERE = Path.cwd()

GENERADOR = str(_HERE)

# Directory containing template files
TEMPLATE_DIR = os.path.join(GENERADOR, "plantillas")

# Input spreadsheet (xlsx or csv)
EXCEL_PATH = os.path.join(GENERADOR, "datos.xlsx")

# Output directory for PDFs
OUTPUT_DIR = os.path.join(GENERADOR, "salida")

# Excel sheet (name or index); first sheet = 0 (ignored for CSV)
SHEET_NAME = 0

# Output PDF filename pattern (can use Excel columns and 'index')
# Examples: "{ID}_{Nombre}.pdf" | "{index:04d}_{Empresa}.pdf"
FILENAME_PATTERN = "{NOMBRE} - {SALIDA}.pdf"

# Token delimiters in templates ({{FIELD}})
TOKEN_PREFIX = "{{"
TOKEN_SUFFIX = "}}"

# PPTX: scan masters/layouts too?
SCAN_MASTERS_PPTX = True

# DOCX: scan headers/footers too?
SCAN_HEADERS_FOOTERS_DOCX = True

# Path to 'soffice' if not in PATH (None to use from PATH)
SOFFICE_BIN = None  # e.g., r"C:\Program Files\LibreOffice\program\soffice.exe"

# Default template if 'TEMPLATE' cell is empty (None -> error)
DEFAULT_TEMPLATE = None  # e.g., "base.pptx"

# Minimum required columns in Excel
REQUIRED_COLUMNS = ["TEMPLATE"]

# Simulate only (no PDF generation)
DRY_RUN = False

# Strict mode: abort if some template token has no matching Excel column
STRICT_MODE = False

# PDF export: retry attempts on failure (total tries = retries + 1)
EXPORT_RETRIES = 2

# Export engine: "auto" | "libreoffice" | "msoffice"
# - auto: prefer msoffice on Windows if available for DOCX/PPTX, else libreoffice
EXPORT_ENGINE = "auto"

# LibreOffice PDF filter options (None for defaults).
# Example: PDF_FILTER="pdf:writer_pdf_Export", PDF_FILTER_OPTS="SelectPdfVersion=1;Quality=90"
PDF_FILTER = "pdf"
PDF_FILTER_OPTS = None

# Supported template extensions and their OOXML targets
OOXML_TARGETS = {
    ".docx": ".docx",
    ".pptx": ".pptx",
    ".doc": ".docx",
    ".ppt": ".pptx",
    ".odt": ".docx",
    ".odp": ".pptx",
    ".rtf": ".docx",
}

# Optional column formatters, applied before plain token substitution.
def _fmt_euros(s: str) -> str:
    try:
        v = float(str(s).replace(".", "").replace(",", "."))
        return f"{v:,.2f} €".replace(",", "X").replace(".", ",").replace("X", ".")
    except Exception:
        return str(s)

def _fmt_fecha_dmy(s: str) -> str:
    import datetime as _dt
    s = str(s).strip()
    if not s:
        return ""
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%d-%m-%Y", "%Y/%m/%d"):
        try:
            d = _dt.datetime.strptime(s, fmt).date()
            return d.strftime("%d/%m/%Y")
        except Exception:
            pass
    return s

COLUMN_FORMATTERS = {
    # "Importe": _fmt_euros,
    # "Fecha": _fmt_fecha_dmy,
}

# Cache for legacy conversions to avoid repeated conversions
_CONVERT_CACHE: Dict[str, Path] = {}

# Try tqdm for progress bar
try:
    from tqdm import tqdm  # type: ignore
except Exception:
    tqdm = None  # progress bar optional

# On Windows, try MS Office COM
_WIN = os.name == "nt"
MSWORD_AVAILABLE = False
MSPPT_AVAILABLE = False
_word_app = None
_ppt_app = None


# =========================
# 🔧 UTILS & LOGGING
# =========================
def _log(msg: str) -> None:
    """Print log line (stdout, flushed)."""
    print(msg, flush=True)

def _sanitize_filename(s: str) -> str:
    """Replace filesystem-unsafe characters and trim."""
    bad = '<>:"/\\|?*'
    for ch in bad:
        s = s.replace(ch, "_")
    return s.strip()

def _ensure_outdir(path: Path) -> None:
    """Ensure directory exists."""
    path.mkdir(parents=True, exist_ok=True)

def _token(col: str) -> str:
    """Build token string from a column name."""
    return f"{TOKEN_PREFIX}{col}{TOKEN_SUFFIX}"

def _normalize_headers_and_values(df: pd.DataFrame) -> pd.DataFrame:
    """Trim headers and cell string values. Use DataFrame.map on pandas 2.2+."""
    df.columns = [str(c).strip() for c in df.columns]
    strip_if_str = lambda x: x.strip() if isinstance(x, str) else x
    if hasattr(pd.DataFrame, "map"):  # pandas >= 2.2
        return df.map(strip_if_str)
    return df.applymap(strip_if_str)

def _apply_formatters(row: pd.Series) -> pd.Series:
    """Apply COLUMN_FORMATTERS to row values (string in, string out)."""
    out = row.copy()
    for col, fn in COLUMN_FORMATTERS.items():
        if col in out:
            try:
                out[col] = fn("" if pd.isna(out[col]) else str(out[col]))
            except Exception:
                out[col] = "" if pd.isna(out[col]) else str(out[col])
    return out

def _read_table(path: str, sheet: Any) -> pd.DataFrame:
    """Read CSV/XLSX by extension and return a normalized DataFrame (strings)."""
    p = str(path).lower()
    if p.endswith(".csv"):
        df = pd.read_csv(path, dtype=str).fillna("")
    else:
        df = pd.read_excel(path, sheet_name=sheet, dtype=str).fillna("")
    return _normalize_headers_and_values(df)


# =========================
# 🔧 TOKEN FILTERS / DEFAULTS
# =========================
# {{Field|trim|upper}}, {{Amount|euros}}, {{Date|dmy}}, {{Field?:N/A}}
TOKEN_INNER_RE = re.compile(r"\{\{\s*([^}]+?)\s*\}\}")

def _eval_token(raw_inner: str, rowdict: Dict[str, str]) -> str:
    """
    Evaluate a token with optional default and filters.
    Grammar (simple):
      <expr> := <main> [ '?:' <default> ]
      <main> := <colname> ( '|' <filter> )*
    Examples:
      "Nombre|trim|upper"
      "Importe|euros"
      "Fecha|dmy"
      "Campo?:N/A"
    """
    default = None
    if "?: " in raw_inner or "?:" in raw_inner:
        parts = raw_inner.split("?:", 1)
        raw_inner, default = parts[0].strip(), parts[1].strip()

    pieces = [p.strip() for p in raw_inner.split("|")] if raw_inner else [raw_inner]
    col = pieces[0] if pieces else ""
    filters = pieces[1:] if len(pieces) > 1 else []

    val = rowdict.get(col, "")
    if default is not None and (val is None or str(val).strip() == ""):
        val = default

    registry = {
        "trim": lambda s: s.strip(),
        "upper": lambda s: s.upper(),
        "lower": lambda s: s.lower(),
        "euros": _fmt_euros,
        "dmy": _fmt_fecha_dmy,
    }
    for f in filters:
        fn = registry.get(f)
        if fn:
            val = fn(str(val))
    return str(val)

def _collect_base_name(inner: str) -> str:
    """
    For preflight, extract the base column name from a complex token.
    Example: "Campo|upper" -> "Campo"; "Campo?:N/A" -> "Campo"
    """
    base = inner.split("?:", 1)[0]
    return base.split("|", 1)[0].strip()


# =========================
# 🔧 LIBREOFFICE CONVERSION/EXPORT
# =========================
def _lo_cmd() -> str:
    """Return the soffice command to execute."""
    return SOFFICE_BIN if SOFFICE_BIN else "soffice"

def lo_convert(input_path: str, outdir: str, to_format: str, filter_opts: Optional[str] = None) -> None:
    """
    Convert using LibreOffice to a target format (pdf, pptx, docx, etc).
    Correct argument order: --convert-to ... [--outdir ...] <file>
    On Windows, hide the spawned console window.
    Raises on non-zero exit.
    """
    conv = to_format if not filter_opts else f"{to_format}:{filter_opts}"
    cmd = [
        _lo_cmd(),
        "--headless",
        "--convert-to", conv,   # must come before --outdir
        "--outdir", outdir,     # outdir must directly follow --convert-to
        input_path,
    ]

    # Hide console window on Windows
    creationflags = 0
    startupinfo = None
    if os.name == "nt":
        try:
            creationflags = subprocess.CREATE_NO_WINDOW  # type: ignore[attr-defined]
        except Exception:
            pass

    proc = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        creationflags=creationflags,
        startupinfo=startupinfo,
    )
    if proc.returncode != 0:
        raise RuntimeError(
            "LibreOffice returned a non-zero exit code.\n"
            f"CMD: {' '.join(cmd)}\n"
            f"STDOUT:\n{proc.stdout}\n\nSTDERR:\n{proc.stderr}"
        )


def export_pdf_with_libreoffice(input_file: str, pdf_out: str) -> None:
    """
    Export to PDF using LibreOffice from the given (pptx/docx) file.
    """
    out_dir = Path(pdf_out).parent.resolve()
    _ensure_outdir(out_dir)

    with tempfile.TemporaryDirectory() as tmpd:
        lo_convert(input_file, tmpd, PDF_FILTER, PDF_FILTER_OPTS)
        produced = Path(tmpd) / (Path(input_file).stem + ".pdf")
        if not produced.exists():
            raise FileNotFoundError("No PDF produced by LibreOffice.")
        shutil.move(str(produced), str(pdf_out))

def export_pdf_with_retry(input_file: str, pdf_out: str, retries: int = EXPORT_RETRIES) -> None:
    """
    Wrapper with simple retries around export_pdf_with_libreoffice.
    """
    last_err = None
    for attempt in range(1, retries + 2):
        try:
            export_pdf_with_libreoffice(input_file, pdf_out)
            return
        except Exception as e:
            last_err = e
            _log(f"[WARN] PDF export attempt {attempt} failed: {e}")
    raise last_err


# =========================
# 🔧 MS OFFICE (COM) EXPORT (Windows only)
# =========================
def _try_init_msoffice() -> None:
    """Attempt to initialize Word and PowerPoint COM on Windows."""
    global MSWORD_AVAILABLE, MSPPT_AVAILABLE, _word_app, _ppt_app
    if not _WIN:
        return
    # Prefer win32com; fallback to comtypes
    try:
        import win32com.client  # type: ignore
        # Word
        try:
            _word_app = win32com.client.Dispatch("Word.Application")
            _word_app.Visible = False
            MSWORD_AVAILABLE = True
        except Exception:
            MSWORD_AVAILABLE = False
        # PowerPoint
        try:
            _ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            _ppt_app.Visible = False
            MSPPT_AVAILABLE = True
        except Exception:
            MSPPT_AVAILABLE = False
        return
    except Exception:
        pass
    # comtypes fallback
    try:
        import comtypes.client  # type: ignore
        try:
            _word_app = comtypes.client.CreateObject("Word.Application")
            _word_app.Visible = False
            MSWORD_AVAILABLE = True
        except Exception:
            MSWORD_AVAILABLE = False
        try:
            _ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
            _ppt_app.Visible = False
            MSPPT_AVAILABLE = True
        except Exception:
            MSPPT_AVAILABLE = False
    except Exception:
        MSWORD_AVAILABLE = False
        MSPPT_AVAILABLE = False

def _msoffice_export_pdf(input_file: str, pdf_out: str) -> bool:
    """
    Try MS Office COM export. Returns True if exported, False if not supported/available.
    - DOCX via Word (wdFormatPDF = 17)
    - PPTX via PowerPoint (ppSaveAsPDF = 32)
    """
    if not _WIN:
        return False
    ext = Path(input_file).suffix.lower()
    out_dir = Path(pdf_out).parent
    _ensure_outdir(out_dir)

    if ext == ".docx" and MSWORD_AVAILABLE and _word_app is not None:
        try:
            doc = _word_app.Documents.Open(str(Path(input_file).resolve()))
            wdFormatPDF = 17
            doc.SaveAs(str(Path(pdf_out).resolve()), FileFormat=wdFormatPDF)
            doc.Close(False)
            return True
        except Exception as e:
            _log(f"[WARN] Word COM export failed: {e}")
            return False

    if ext == ".pptx" and MSPPT_AVAILABLE and _ppt_app is not None:
        try:
            pres = _ppt_app.Presentations.Open(str(Path(input_file).resolve()), WithWindow=False)
            ppSaveAsPDF = 32
            pres.SaveAs(str(Path(pdf_out).resolve()), ppSaveAsPDF)
            pres.Close()
            return True
        except Exception as e:
            _log(f"[WARN] PowerPoint COM export failed: {e}")
            return False

    return False


# =========================
# 🔧 CONVERSIONS TO OOXML
# =========================
def convert_to_ooxml_if_needed(template_path: Path) -> Path:
    """
    Convert templates (doc, ppt, odt, odp, rtf) to OOXML (docx/pptx) once and cache the result.
    Returns the OOXML path (original path if already OOXML).
    """
    ext = template_path.suffix.lower()
    if ext in (".docx", ".pptx"):
        return template_path
    if ext not in OOXML_TARGETS:
        raise ValueError(f"Unsupported template extension: {ext}")

    key = str(template_path.resolve())
    if key in _CONVERT_CACHE:
        return _CONVERT_CACHE[key]

    target_ext = OOXML_TARGETS[ext]
    tmpd = tempfile.mkdtemp()
    lo_convert(str(template_path), tmpd, target_ext.lstrip("."))  # "docx" or "pptx"
    converted = Path(tmpd) / (template_path.stem + target_ext)

    if not converted.exists():
        raise FileNotFoundError(f"Conversion to OOXML failed: {converted}")

    _CONVERT_CACHE[key] = converted
    return converted


# =========================
# 🔧 PPTX TOKEN REPLACEMENT
# =========================
def _pptx_replace_in_paragraph(paragraph, mapping: dict, rowdict: Dict[str, str]) -> None:
    """
    Replace tokens within a PPTX paragraph.
    - First, fast replace for simple mapping tokens {{Col}}.
    - Then, evaluate complex tokens {{Col|filter|...}} and {{Col?:default}}.
    - Collapses paragraph to a single run (preserving the first run's basic formatting).
    """
    full_text = "".join(run.text for run in paragraph.runs)
    if not full_text:
        return

    replaced = full_text
    changed = False

    # 1) Fast path for plain mapping tokens
    for token, value in mapping.items():
        if token in replaced:
            replaced = replaced.replace(token, value)
            changed = True

    # 2) Complex tokens evaluation
    def _eval_match(m: re.Match) -> str:
        inner = m.group(1)
        return _eval_token(inner, rowdict)

    if TOKEN_INNER_RE.search(replaced):
        new_text = TOKEN_INNER_RE.sub(_eval_match, replaced)
        if new_text != replaced:
            replaced = new_text
            changed = True

    if not changed:
        return

    # Preserve basic formatting from the first run
    if paragraph.runs:
        fr = paragraph.runs[0]
        base = {
            "size": fr.font.size,
            "bold": fr.font.bold,
            "italic": fr.font.italic,
            "name": fr.font.name,
            "underline": fr.font.underline,
            "color_rgb": getattr(fr.font.color, "rgb", None),
        }
    else:
        base = dict(size=None, bold=None, italic=None, name=None, underline=None, color_rgb=None)

    for _ in list(paragraph.runs):
        paragraph._p.remove(paragraph.runs[0]._r)

    new_run = paragraph.add_run()
    new_run.text = replaced
    if base["size"] is not None: new_run.font.size = base["size"]
    if base["bold"] is not None: new_run.font.bold = base["bold"]
    if base["italic"] is not None: new_run.font.italic = base["italic"]
    if base["name"] is not None: new_run.font.name = base["name"]
    if base["underline"] is not None: new_run.font.underline = base["underline"]
    if base["color_rgb"] is not None: new_run.font.color.rgb = base["color_rgb"]

def _pptx_iter_table_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def _pptx_walk_shapes(shapes, mapping: dict, rowdict: Dict[str, str]) -> None:
    for shape in shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for p in shape.text_frame.paragraphs:
                _pptx_replace_in_paragraph(p, mapping, rowdict)
        if getattr(shape, "has_table", False):
            for cell in _pptx_iter_table_cells(shape.table):
                if cell.text_frame:
                    for p in cell.text_frame.paragraphs:
                        _pptx_replace_in_paragraph(p, mapping, rowdict)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            _pptx_walk_shapes(shape.shapes, mapping, rowdict)

def _pptx_replace_on_slides(prs: Presentation, mapping: dict, rowdict: Dict[str, str]) -> None:
    for slide in prs.slides:
        _pptx_walk_shapes(slide.shapes, mapping, rowdict)
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                for p in slide.notes_slide.notes_text_frame.paragraphs:
                    _pptx_replace_in_paragraph(p, mapping, rowdict)
        except Exception:
            pass

def _pptx_replace_on_masters(prs: Presentation, mapping: dict, rowdict: Dict[str, str]) -> None:
    try:
        for master in prs.slide_masters:
            _pptx_walk_shapes(master.shapes, mapping, rowdict)
            for layout in master.slide_layouts:
                _pptx_walk_shapes(layout.shapes, mapping, rowdict)
    except Exception:
        pass

def replace_placeholders_pptx(prs: Presentation, mapping: dict, rowdict: Dict[str, str], scan_masters: bool) -> None:
    if scan_masters:
        _pptx_replace_on_masters(prs, mapping, rowdict)
    _pptx_replace_on_slides(prs, mapping, rowdict)


# =========================
# 🔧 DOCX TOKEN REPLACEMENT
# =========================
def _docx_replace_in_paragraph(paragraph, mapping: dict, rowdict: Dict[str, str]) -> None:
    """
    Replace tokens in a DOCX paragraph safely.
    - Fast replace for mapping tokens {{Col}}.
    - Evaluate complex tokens {{Col|...}} and {{Col?:...}}.
    - Build the final text and put it into the first run.
    - Remove remaining runs using a snapshot of XML elements to avoid
      'Element is not a child of this node'.
    - Keep basic run formatting and paragraph style from the first run.
    """
    runs = paragraph.runs
    full_text = "".join(r.text for r in runs)
    if not full_text:
        return

    replaced = full_text
    changed = False

    for token, value in mapping.items():
        if token in replaced:
            replaced = replaced.replace(token, value)
            changed = True

    def _eval_match(m: re.Match) -> str:
        inner = m.group(1)
        return _eval_token(inner, rowdict)

    if TOKEN_INNER_RE.search(replaced):
        new_text = TOKEN_INNER_RE.sub(_eval_match, replaced)
        if new_text != replaced:
            replaced = new_text
            changed = True

    if not changed:
        return

    if runs:
        fr = runs[0]
        base_bold = fr.bold
        base_italic = fr.italic
        base_underline = fr.underline
        para_style_name = paragraph.style.name if paragraph.style else None
    else:
        base_bold = base_italic = base_underline = None
        para_style_name = paragraph.style.name if paragraph.style else None

    # Snapshot original XML run elements before modifying
    r_elems = [r._r for r in runs]

    # Ensure at least one run exists
    if not runs:
        paragraph.add_run("")
        runs = paragraph.runs
        r_elems = [r._r for r in runs]

    # Write all text into the first run
    first_run = paragraph.runs[0]
    first_run.text = replaced

    # Restore basic formatting and paragraph style
    if base_bold is not None: first_run.bold = base_bold
    if base_italic is not None: first_run.italic = base_italic
    if base_underline is not None: first_run.underline = base_underline
    if para_style_name:
        try:
            paragraph.style = para_style_name
        except Exception:
            pass

    # Remove remaining runs using the XML snapshot
    for r in r_elems[1:]:
        parent = r.getparent()
        if parent is not None:
            try:
                parent.remove(r)
            except Exception:
                # If already removed/reparented (fields, hyperlinks), ignore
                pass

def _docx_walk_table(table, mapping: dict, rowdict: Dict[str, str]) -> None:
    for row in table.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                _docx_replace_in_paragraph(p, mapping, rowdict)
            for t in cell.tables:
                _docx_walk_table(t, mapping, rowdict)

def replace_placeholders_docx(doc: DocxDocument, mapping: dict, rowdict: Dict[str, str], scan_headers_footers: bool) -> None:
    # Body
    for p in doc.paragraphs:
        _docx_replace_in_paragraph(p, mapping, rowdict)
    for t in doc.tables:
        _docx_walk_table(t, mapping, rowdict)

    # Headers/Footers
    if scan_headers_footers:
        for section in doc.sections:
            header = section.header
            if header:
                for p in header.paragraphs:
                    _docx_replace_in_paragraph(p, mapping, rowdict)
                for t in header.tables:
                    _docx_walk_table(t, mapping, rowdict)
            footer = section.footer
            if footer:
                for p in footer.paragraphs:
                    _docx_replace_in_paragraph(p, mapping, rowdict)
                for t in footer.tables:
                    _docx_walk_table(t, mapping, rowdict)


# =========================
# 🔧 TOKEN DISCOVERY (PREFLIGHT)
# =========================
# Collect inner content; later extract base names for column check
TOKEN_RE = re.compile(r"\{\{\s*([A-Za-z0-9_\- :|]+)\s*\}\}")

def _pptx_collect_tokens(prs: Presentation) -> set[str]:
    """Collect all {{token}} inner strings from a PPTX Presentation."""
    found: set[str] = set()

    def _collect_paragraph(p):
        txt = "".join(r.text for r in p.runs) if getattr(p, "runs", None) else (p.text or "")
        for m in TOKEN_RE.finditer(txt):
            found.add(m.group(1).strip())

    def _walk(shapes):
        for s in shapes:
            if hasattr(s, "text_frame") and s.text_frame:
                for p in s.text_frame.paragraphs:
                    _collect_paragraph(p)
            if getattr(s, "has_table", False):
                for row in s.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for p in cell.text_frame.paragraphs:
                                _collect_paragraph(p)
            if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(s, "shapes"):
                _walk(s.shapes)

    if SCAN_MASTERS_PPTX:
        try:
            for master in prs.slide_masters:
                _walk(master.shapes)
                for layout in master.slide_layouts:
                    _walk(layout.shapes)
        except Exception:
            pass

    for slide in prs.slides:
        _walk(slide.shapes)
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                for p in slide.notes_slide.notes_text_frame.paragraphs:
                    _collect_paragraph(p)
        except Exception:
            pass

    return found

def _docx_collect_tokens(doc: DocxDocument) -> set[str]:
    """Collect all {{token}} inner strings from a DOCX Document."""
    found: set[str] = set()

    def _collect_paragraph(p):
        txt = "".join(r.text for r in p.runs) if getattr(p, "runs", None) else (p.text or "")
        for m in TOKEN_RE.finditer(txt):
            found.add(m.group(1).strip())

    def _walk_table(table):
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _collect_paragraph(p)
                for t in cell.tables:
                    _walk_table(t)

    for p in doc.paragraphs:
        _collect_paragraph(p)
    for t in doc.tables:
        _walk_table(t)

    if SCAN_HEADERS_FOOTERS_DOCX:
        for section in doc.sections:
            h = section.header
            if h:
                for p in h.paragraphs: _collect_paragraph(p)
                for t in h.tables: _walk_table(t)
            f = section.footer
            if f:
                for p in f.paragraphs: _collect_paragraph(p)
                for t in f.tables: _walk_table(t)

    return found

def _collect_tokens_from_template_path(template_path: Path) -> set[str]:
    """Convert to OOXML if needed and collect tokens from the OOXML file."""
    ooxml = convert_to_ooxml_if_needed(template_path)
    if ooxml.suffix.lower() == ".pptx":
        prs = Presentation(str(ooxml))
        return _pptx_collect_tokens(prs)
    elif ooxml.suffix.lower() == ".docx":
        doc = DocxDocument(str(ooxml))
        return _docx_collect_tokens(doc)
    return set()


# =========================
# 🔧 TEMPLATE RESOLUTION / DISPATCH
# =========================
def _resolve_template_path(filename: str) -> Path:
    """
    Validate 'TEMPLATE' (filename + extension, no directories),
    and return TEMPLATE_DIR/<filename>. If empty, use DEFAULT_TEMPLATE if provided.
    """
    name = (filename or "").strip()
    if not name:
        if DEFAULT_TEMPLATE:
            name = DEFAULT_TEMPLATE
        else:
            raise ValueError("Excel column 'TEMPLATE' is empty and no DEFAULT_TEMPLATE is configured.")
    if any(ch in name for ch in ("/", "\\")):
        raise ValueError(f"'TEMPLATE' must be a filename only (no directories). Received: {name!r}")
    path = Path(TEMPLATE_DIR) / name
    if not path.exists():
        raise FileNotFoundError(f"Template file not found: {path}")
    return path

def _export_pdf_with_selected_engine(edited_path: Path, pdf_path: Path) -> None:
    """
    Export PDF using selected engine:
      - msoffice (Word/PowerPoint COM on Windows) for DOCX/PPTX if available,
      - otherwise LibreOffice.
    """
    engine = EXPORT_ENGINE.lower().strip()
    ext = edited_path.suffix.lower()

    # auto → prefer MS Office on Windows if app is ready and file is OOXML
    if engine == "auto":
        if _WIN and ((ext == ".docx" and MSWORD_AVAILABLE) or (ext == ".pptx" and MSPPT_AVAILABLE)):
            if _msoffice_export_pdf(str(edited_path), str(pdf_path)):
                return
        # fallback to LibreOffice
        export_pdf_with_retry(str(edited_path), str(pdf_path))
        return

    if engine == "msoffice":
        # Only for OOXML; else we fallback to LO
        if _WIN and ext in (".docx", ".pptx"):
            if _msoffice_export_pdf(str(edited_path), str(pdf_path)):
                return
        # Fallback
        export_pdf_with_retry(str(edited_path), str(pdf_path))
        return

    # Default: LibreOffice
    export_pdf_with_retry(str(edited_path), str(pdf_path))

def render_pdf_from_template(template_path: Path, rowdict: Dict[str, str], pdf_path: Path) -> None:
    """
    Dispatch by extension:
      - Convert to OOXML if needed (.doc/.ppt/.odt/.odp/.rtf)
      - Replace tokens
      - Export PDF using selected engine
    """
    ext = template_path.suffix.lower()
    if ext not in OOXML_TARGETS:
        # still allow direct OOXML
        if ext not in (".docx", ".pptx"):
            raise ValueError(f"Unsupported extension: {ext}")

    ooxml_path = convert_to_ooxml_if_needed(template_path) if ext != ".docx" and ext != ".pptx" else template_path
    ooxml_ext = ooxml_path.suffix.lower()

    # Build mapping for fast replacement of simple tokens
    # (Row already normalized; apply formatters)
    mapping = {}
    for k, v in rowdict.items():
        if k.upper() == "TEMPLATE":
            continue
        mapping[_token(k)] = str(v)

    with tempfile.TemporaryDirectory() as tmpd:
        tmpd_path = Path(tmpd)

        if ooxml_ext == ".pptx":
            prs = Presentation(str(ooxml_path))
            replace_placeholders_pptx(prs, mapping, rowdict, scan_masters=SCAN_MASTERS_PPTX)
            edited = tmpd_path / "edited.pptx"
            prs.save(edited)
            _export_pdf_with_selected_engine(edited, pdf_path)

        elif ooxml_ext == ".docx":
            doc = DocxDocument(str(ooxml_path))
            replace_placeholders_docx(doc, mapping, rowdict, scan_headers_footers=SCAN_HEADERS_FOOTERS_DOCX)
            edited = tmpd_path / "edited.docx"
            doc.save(edited)
            _export_pdf_with_selected_engine(edited, pdf_path)

        else:
            raise RuntimeError(f"Unexpected OOXML extension after conversion: {ooxml_ext}")


# =========================
# 🔧 BATCH PROCESS (SPREADSHEET-DRIVEN)
# =========================
def run_batch() -> None:
    # Use pre-loaded DataFrame if present (testing/filters), else read from disk
    df = globals().get("_DF_OVERRIDE")
    if df is None:
        df = _read_table(EXCEL_PATH, SHEET_NAME)

    # Validate minimum required columns
    missing = [c for c in REQUIRED_COLUMNS if c not in df.columns]
    if missing:
        raise KeyError(f"Missing required Excel columns: {missing}. Available: {list(df.columns)}")

    out_dir = Path(OUTPUT_DIR)
    _ensure_outdir(out_dir)

    # --- Preflight: tokens vs. Excel columns ---
    excel_cols = set(c for c in df.columns if c.upper() != "TEMPLATE")
    used_templates = set(df["TEMPLATE"].dropna().astype(str).str.strip()) if "TEMPLATE" in df.columns else set()
    if not used_templates and DEFAULT_TEMPLATE:
        used_templates = {DEFAULT_TEMPLATE}

    all_tokens_raw = set()
    for tplname in used_templates:
        tpath = _resolve_template_path(tplname)
        toks = _collect_tokens_from_template_path(tpath)
        all_tokens_raw |= toks

    base_tokens = {_collect_base_name(t) for t in all_tokens_raw}
    missing_cols = {t for t in base_tokens if t not in excel_cols}
    unused_cols = {c for c in excel_cols if c not in base_tokens}

    _log(f"[Preflight] Tokens found in templates (raw): {sorted(all_tokens_raw)}")
    _log(f"[Preflight] Base token names: {sorted(base_tokens)}")
    if missing_cols:
        _log(f"[Preflight][WARN] Tokens without matching Excel columns: {sorted(missing_cols)}")
    if unused_cols:
        _log(f"[Preflight][INFO] Excel columns not used by any token: {sorted(unused_cols)}")
    if STRICT_MODE and missing_cols:
        raise KeyError(f"[STRICT] Missing Excel columns for tokens: {sorted(missing_cols)}")

    total = len(df)
    _log(f"Rows to process: {total}")

    results = []
    iterator = df.iterrows()
    if tqdm is not None:
        iterator = tqdm(df.iterrows(), total=total, desc="Processing", unit="row")

    for i, row in iterator:
        # Row-level skip (column 'SKIP')
        skip_val = str(row.get("SKIP", "")).strip().lower()
        if skip_val in ("1", "true", "sí", "si", "x", "y", "yes"):
            _log(f"[{i+1}/{total}] SKIP=1 → row skipped")
            results.append({"row": int(i), "status": "SKIPPED"})
            continue

        template_name = row.get("TEMPLATE", "")
        try:
            template_path = _resolve_template_path(template_name)
        except Exception as e:
            _log(f"[{i+1}/{total}] [ERROR] Template resolve failed: {e}")
            results.append({"row": int(i), "status": "ERROR", "error": str(e)})
            continue

        # Build row dict (stringified) with formatters applied for simple mapping
        row2 = row.fillna("").astype(str)
        row2 = _apply_formatters(row2)
        rowdict = row2.to_dict()

        # Output filename
        ctx = {**rowdict, "index": i}
        try:
            pdf_name = FILENAME_PATTERN.format(**ctx)
        except KeyError as e:
            err = (f"Filename pattern requires a missing column: {e}. "
                   f"Pattern: {FILENAME_PATTERN} | Columns: {list(df.columns)}")
            _log(f"[{i+1}/{total}] [ERROR] {err}")
            results.append({"row": int(i), "status": "ERROR", "error": err})
            continue

        pdf_name = _sanitize_filename(pdf_name)
        if not pdf_name.lower().endswith(".pdf"):
            pdf_name += ".pdf"

        # Optional per-row subfolder (column 'OUTPUT')
        subdir = str(row.get("OUTPUT", "")).strip()
        target_dir = Path(OUTPUT_DIR) / _sanitize_filename(subdir) if subdir else Path(OUTPUT_DIR)
        _ensure_outdir(target_dir)
        pdf_path = target_dir / pdf_name

        # Dry-run
        if DRY_RUN:
            _log(f"[{i+1}/{total}] [DRY-RUN] Template={template_path.name} -> {pdf_path.name}")
            results.append({
                "row": int(i),
                "status": "DRY-RUN",
                "template": template_name,
                "output": str(pdf_path),
            })
            continue

        # Generate
        try:
            _log(f"[{i+1}/{total}] {template_path.name} → {pdf_path.name}")
            render_pdf_from_template(template_path, rowdict, pdf_path)
            size = Path(pdf_path).stat().st_size if Path(pdf_path).exists() else 0
            results.append({
                "row": int(i),
                "status": "OK",
                "template": template_name,
                "output": str(pdf_path),
                "bytes": int(size),
            })
        except Exception as e:
            results.append({
                "row": int(i),
                "status": "ERROR",
                "template": template_name,
                "output": str(pdf_path),
                "error": str(e),
            })
            _log(f"[ERROR] Row {i} ({template_path.name}) → {e}")
            continue

    # Save reports
    try:
        rep_path = Path(OUTPUT_DIR) / "_report.json"
        with open(rep_path, "w", encoding="utf-8") as f:
            json.dump(results, f, ensure_ascii=False, indent=2)
        _log(f"Report saved to: {rep_path}")
    except Exception as e:
        _log(f"[WARN] Could not write JSON report: {e}")

    try:
        csv_path = Path(OUTPUT_DIR) / "_report.csv"
        keys = sorted({k for r in results for k in r.keys()})
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            w = csv.DictWriter(f, fieldnames=keys)
            w.writeheader()
            w.writerows(results)
        _log(f"Report saved to: {csv_path}")
    except Exception as e:
        _log(f"[WARN] Could not write CSV report: {e}")

    _log("✅ Done.")


# =========================
# 🔧 MAIN (CLI OVERRIDES)
# =========================

def _setup_logging(verbose: bool) -> None:
    lvl = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(level=lvl, format="%(message)s")
    # Redirect our _log to logging.info for consistency
    global _log
    _log = lambda msg: logging.info(msg)

def _detect_soffice() -> tuple[bool, str]:
    exe = SOFFICE_BIN or "soffice"
    try:
        proc = subprocess.run([exe, "--version"], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        ok = proc.returncode == 0
        out = proc.stdout.strip() or proc.stderr.strip()
        return ok, out
    except Exception as e:
        return False, str(e)

def _detect_msoffice() -> dict:
    return {
        "windows": _WIN,
        "word_available": MSWORD_AVAILABLE,
        "powerpoint_available": MSPPT_AVAILABLE,
    }

def _apply_row_filters(df: pd.DataFrame, r_from: int | None, r_to: int | None, where: str | None) -> pd.DataFrame:
    if where:
        try:
            df = df.query(where, engine="python")  # simple filters
        except Exception as e:
            _log(f"[WARN] --where ignored due to error: {e}")
    if r_from is not None or r_to is not None:
        start = 0 if r_from is None else int(r_from)
        stop  = len(df) if r_to is None else int(r_to) + 1
        df = df.iloc[start:stop]
    return df

def run_batch_from_dataframe(df: pd.DataFrame) -> None:
    """
    Same as run_batch(), but operates on an already-loaded DataFrame.
    Keeps behavior for reporting and generation. This makes testing easier.
    """
    # The body is identical to run_batch() up to `df = _read_table(...)`.
    # To avoid duplication, we can lift the core loop into a helper and call it here.
    # For now, reuse the existing run_batch() implementation path by temporarily
    # monkey-patching a global (minimal intrusion).
    global _DF_OVERRIDE
    _DF_OVERRIDE = df
    try:
        run_batch()
    finally:
        _DF_OVERRIDE = None

def main():

    # Override globals from CLI (non-breaking)
    global EXCEL_PATH, OUTPUT_DIR, TEMPLATE_DIR, SHEET_NAME, FILENAME_PATTERN
    global EXPORT_ENGINE, DRY_RUN, STRICT_MODE, PDF_FILTER_OPTS

    parser = argparse.ArgumentParser(description="Batch Office templating ({{TOKENS}}) to PDF.")
    parser.add_argument("data", nargs="?", default=EXCEL_PATH, help="Input XLSX/CSV path (default from config).")
    parser.add_argument("outdir", nargs="?", default=OUTPUT_DIR, help="Output directory (default from config).")
    parser.add_argument("templates", nargs="?", default=TEMPLATE_DIR, help="Templates directory (default from config).")
    parser.add_argument("--sheet", default=SHEET_NAME, help="Excel sheet name or index (ignored for CSV).")
    parser.add_argument("--pattern", default=FILENAME_PATTERN, help="Output filename pattern.")
    parser.add_argument("--engine", choices=["auto", "libreoffice", "msoffice"], default=EXPORT_ENGINE)
    parser.add_argument("--strict", action="store_true", default=STRICT_MODE, help="Fail if any token has no column.")
    parser.add_argument("--dry-run", action="store_true", default=DRY_RUN, help="Do not generate PDFs.")
    parser.add_argument("--pdf-filter-opts", default=PDF_FILTER_OPTS, help="LibreOffice PDF filter options.")
    parser.add_argument("--from", dest="r_from", type=int, help="Start row (inclusive).")
    parser.add_argument("--to", dest="r_to", type=int, help="End row (inclusive).")
    parser.add_argument("--where", help='Pandas query filter, e.g. "Curso == \'A\' and SKIP != \'1\'"')
    parser.add_argument("--verbose", action="store_true", help="Verbose logging.")
    parser.add_argument("--version", action="store_true", help="Print version and exit.")
    parser.add_argument("--check", action="store_true", help="Check environment and exit.")
    args = parser.parse_args()

    if args.version:
        print(__version__)
        return

    _setup_logging(args.verbose)

    EXCEL_PATH = args.data
    OUTPUT_DIR = args.outdir
    TEMPLATE_DIR = args.templates
    SHEET_NAME = args.sheet
    FILENAME_PATTERN = args.pattern
    EXPORT_ENGINE = args.engine
    DRY_RUN = bool(args.dry_run)
    STRICT_MODE = bool(args.strict)
    PDF_FILTER_OPTS = args.pdf_filter_opts

    # Environment check
    if args.check:
        ok_soff, info = _detect_soffice()
        _log(f"[Check] LibreOffice: {'OK' if ok_soff else 'NOT FOUND'} ({info})")
        if _WIN:
            _try_init_msoffice()
            ms = _detect_msoffice()
            _log(f"[Check] Windows: {ms['windows']} | Word: {ms['word_available']} | PowerPoint: {ms['powerpoint_available']}")
        return

    # MS Office init if needed
    if EXPORT_ENGINE in ("auto", "msoffice") and _WIN:
        _try_init_msoffice()

        # Ensure COM apps quit at process end
        def _teardown_office():
            try:
                if _ppt_app is not None:
                    _ppt_app.Quit()
            except Exception:
                pass
            try:
                if _word_app is not None:
                    _word_app.Quit()
            except Exception:
                pass
        atexit.register(_teardown_office)

    # Read and filter
    df = _read_table(EXCEL_PATH, SHEET_NAME)
    df = _apply_row_filters(df, args.r_from, args.r_to, args.where)

    # Run
    run_batch_from_dataframe(df)

if __name__ == "__main__":
    main()
